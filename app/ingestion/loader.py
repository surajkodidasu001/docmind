"""Multi-format document loader.

Supports: pdf, docx, pptx, txt, md, csv, html, json (8 formats).
Each loader returns plain text plus lightweight per-page/section markers
so the chunker can preserve provenance for citations.
"""
from __future__ import annotations

import json
import csv
import io
from dataclasses import dataclass
from pathlib import Path
from typing import List


@dataclass
class RawSegment:
    text: str
    source: str          # file name
    location: str        # page/slide/row label for citation


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md", ".csv", ".html", ".htm", ".json"}


def load_document(path: str, vision_fallback: bool = False) -> List[RawSegment]:
    p = Path(path)
    ext = p.suffix.lower()
    name = p.name

    if ext == ".pdf":
        return _load_pdf(p, name, vision_fallback=vision_fallback)
    if ext == ".docx":
        return _load_docx(p, name)
    if ext == ".pptx":
        return _load_pptx(p, name)
    if ext in (".txt", ".md"):
        return _load_text(p, name)
    if ext == ".csv":
        return _load_csv(p, name)
    if ext in (".html", ".htm"):
        return _load_html(p, name)
    if ext == ".json":
        return _load_json(p, name)

    raise ValueError(f"Unsupported file type: {ext}. Supported: {sorted(SUPPORTED_EXTENSIONS)}")


def _load_pdf(p: Path, name: str, vision_fallback: bool = False) -> List[RawSegment]:
    if vision_fallback:
        # opt-in: costs one vision API call per image-only page, so this
        # isn't the default path — see app/ingestion/vision_pdf.py
        from app.ingestion.vision_pdf import extract_pdf_with_vision_fallback
        page_texts = extract_pdf_with_vision_fallback(str(p))
        return [
            RawSegment(text=text, source=name, location=f"page {i}")
            for i, text in enumerate(page_texts, start=1)
            if text.strip()
        ]

    from pypdf import PdfReader

    reader = PdfReader(str(p))
    segments = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            segments.append(RawSegment(text=text, source=name, location=f"page {i}"))
    return segments


def _load_docx(p: Path, name: str) -> List[RawSegment]:
    import docx

    d = docx.Document(str(p))
    segments = []
    buffer = []
    section = 1
    for para in d.paragraphs:
        if para.text.strip():
            buffer.append(para.text)
        if len("\n".join(buffer)) > 1500:
            segments.append(RawSegment(text="\n".join(buffer), source=name, location=f"section {section}"))
            buffer = []
            section += 1
    if buffer:
        segments.append(RawSegment(text="\n".join(buffer), source=name, location=f"section {section}"))
    return segments


def _load_pptx(p: Path, name: str) -> List[RawSegment]:
    from pptx import Presentation

    prs = Presentation(str(p))
    segments = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        joined = "\n".join(t for t in texts if t.strip())
        if joined.strip():
            segments.append(RawSegment(text=joined, source=name, location=f"slide {i}"))
    return segments


def _load_text(p: Path, name: str) -> List[RawSegment]:
    text = p.read_text(encoding="utf-8", errors="ignore")
    return [RawSegment(text=text, source=name, location="document")]


def _load_csv(p: Path, name: str) -> List[RawSegment]:
    """Chunks CSV rows into text for semantic retrieval (existing behavior),
    and separately stores the parsed DataFrame in the tabular store so
    aggregation-style questions can be answered with real pandas computation
    instead of a text approximation. See app/agent/tabular_compute.py."""
    import pandas as pd
    from app.retrieval.tabular_store import store_dataframe

    df = pd.read_csv(p)
    store_dataframe(name, df)

    segments = []
    with open(p, newline="", encoding="utf-8", errors="ignore") as f:
        reader = csv.DictReader(f)
        for i, row in enumerate(reader, start=1):
            row_text = ", ".join(f"{k}: {v}" for k, v in row.items())
            segments.append(RawSegment(text=row_text, source=name, location=f"row {i}"))
            if i >= 500:  # safety cap for MVP
                break
    return segments


def _load_html(p: Path, name: str) -> List[RawSegment]:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(p.read_text(encoding="utf-8", errors="ignore"), "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    text = soup.get_text(separator="\n")
    text = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    return [RawSegment(text=text, source=name, location="document")]


def _load_json(p: Path, name: str) -> List[RawSegment]:
    data = json.loads(p.read_text(encoding="utf-8", errors="ignore"))
    text = json.dumps(data, indent=2)
    return [RawSegment(text=text, source=name, location="document")]
